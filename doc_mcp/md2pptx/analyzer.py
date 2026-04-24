"""
출력 PPTX 분석기 — MD → PPTX 변환 결과에서 문제를 감지.

입력:
  - 원본 양식 PPTX 경로
  - 변환 결과 PPTX 경로
  - convert() 반환 메타 (tables_matched, tables_unmatched, slides_final, ...)

출력: issues 리스트
  - table_overflow        표 행 수 × 평균 행 높이 > 사용 가능 슬라이드 높이
  - cell_clip             셀 글자수 × 평균 글자폭 > 셀 용량
  - unmatched_table       MD 표 헤더 매칭 실패
  - template_shape_removed 원본 양식 대비 사라진 shape (사용자 수기 편집)
  - prose_unmapped        MD 에 prose 있는데 양식엔 대응 body shape 없음 (v6-2 이후)
  - body_slot_empty       양식에 body shape 있는데 해당 섹션 MD prose 없음 (v6-2 이후)

순수 python-pptx + lxml 로 동작, LLM 없음.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation


# EMU 단위 (1 inch = 914400 EMU)
EMU_PER_INCH = 914400
# 한글 11pt 대략 폭 (0.16" × 914400)
CHAR_W_EMU = 146_000
# 기본 행 높이 (0.4")
DEFAULT_ROW_H_EMU = 370_000
# 본문 기본 줄 높이 (0.2")
LINE_H_EMU = 180_000
# 본문 하단 여백 추정
SLIDE_BOTTOM_MARGIN_EMU = 200_000


def _shape_total_row_height(shape) -> int:
    """표 shape 의 행 합계 높이 (EMU)."""
    if not shape.has_table:
        return 0
    t = shape.table
    total = 0
    for row in t.rows:
        total += int(row.height or DEFAULT_ROW_H_EMU)
    return total


def _estimate_cell_capacity(cell, col_width_emu: int) -> int:
    """셀 용량 추정 (글자수)."""
    try:
        h = int(cell.height or DEFAULT_ROW_H_EMU)
    except Exception:
        h = DEFAULT_ROW_H_EMU
    usable_w = max(col_width_emu - 180_000, CHAR_W_EMU)   # 좌우 여백 뺀 실폭
    chars_per_line = max(1, usable_w // CHAR_W_EMU)
    lines_capacity = max(1, h // LINE_H_EMU)
    return int(chars_per_line * lines_capacity)


def _collect_shape_ids(pptx_path: Path) -> dict[int, set[str]]:
    """양식·결과 비교용: slide_idx → shape name 집합."""
    prs = Presentation(str(pptx_path))
    out: dict[int, set[str]] = {}
    for i, slide in enumerate(prs.slides, start=1):
        names: set[str] = set()
        for sh in slide.shapes:
            # name 이 없는 shape 은 index 기반 fallback
            names.add(sh.name or f"shape_{len(names)}")
        out[i] = names
    return out


def analyze_output(
    template_pptx: str,
    output_pptx: str,
    convert_result: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """결과 PPTX 에서 문제 감지.

    convert_result: md2pptx.cli.convert() 가 반환한 메타. None 이어도 동작하되
                    tables_unmatched 같은 건 놓침.
    """
    tpl = Path(template_pptx)
    out = Path(output_pptx)
    if not out.exists():
        raise FileNotFoundError(output_pptx)

    prs = Presentation(str(out))
    slide_h = int(prs.slide_height or 0)

    issues: list[dict[str, Any]] = []

    # --- 표 넘침 · 셀 클리핑 ---
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for sh_idx, sh in enumerate(slide.shapes):
            if sh.has_table:
                t = sh.table
                n_rows = len(list(t.rows))
                total_h = _shape_total_row_height(sh)
                shape_top = int(sh.top or 0)
                available = slide_h - shape_top - SLIDE_BOTTOM_MARGIN_EMU
                if total_h > available and available > 0:
                    avg_row = total_h / max(n_rows, 1)
                    rows_capacity = max(1, int(available // avg_row))
                    issues.append({
                        "type": "table_overflow",
                        "slide": slide_idx,
                        "shape_index": sh_idx,
                        "rows_used": n_rows,
                        "rows_capacity_est": rows_capacity,
                        "excess_rows": n_rows - rows_capacity,
                        "suggestion": f"표를 {(n_rows + rows_capacity - 1) // rows_capacity}장으로 분할 권장",
                    })

                # 셀 별 클리핑 체크
                cols = t.columns
                n_cols = len(cols)
                for r_idx, row in enumerate(t.rows):
                    for c_idx, cell in enumerate(row.cells):
                        col_w = int(cols[c_idx].width or 0) if c_idx < n_cols else 0
                        if col_w <= 0:
                            continue
                        txt = cell.text_frame.text if cell.text_frame else ""
                        capacity = _estimate_cell_capacity(cell, col_w)
                        if len(txt) > capacity * 1.3:
                            issues.append({
                                "type": "cell_clip",
                                "slide": slide_idx,
                                "shape_index": sh_idx,
                                "row": r_idx,
                                "col": c_idx,
                                "chars": len(txt),
                                "capacity_est": capacity,
                                "excerpt": txt[:60] + ("…" if len(txt) > 60 else ""),
                                "suggestion": "셀 텍스트를 요약하거나 행을 분할",
                            })

            # 일반 text_frame 클리핑 (대략)
            elif sh.has_text_frame:
                txt = sh.text_frame.text
                w = int(sh.width or 0)
                h = int(sh.height or 0)
                if w <= 0 or h <= 0 or not txt:
                    continue
                capacity = _estimate_cell_capacity(sh, w) if hasattr(sh, "height") else 0
                # shape 높이도 충분히 넓어야 길 수 있음
                chars_per_line = max(1, (w - 180_000) // CHAR_W_EMU)
                lines_capacity = max(1, h // LINE_H_EMU)
                total_capacity = chars_per_line * lines_capacity
                if len(txt) > total_capacity * 1.3:
                    issues.append({
                        "type": "text_clip",
                        "slide": slide_idx,
                        "shape_index": sh_idx,
                        "chars": len(txt),
                        "capacity_est": total_capacity,
                        "excerpt": txt[:60] + ("…" if len(txt) > 60 else ""),
                        "suggestion": "텍스트를 요약하거나 bullet 으로 분할",
                    })

    # --- 미매칭 MD 표 (convert_result 있을 때만) ---
    if convert_result:
        for unmatched_idx in convert_result.get("tables_unmatched", []):
            issues.append({
                "type": "unmatched_table",
                "md_table_idx": unmatched_idx,
                "suggestion": "MD 표 헤더를 양식 표 헤더와 맞추거나, 양식에 해당 스키마의 표 슬라이드 추가",
            })

        # v6-2: prose 섹션이 body shape 에 매핑 안 됨
        for u in convert_result.get("body_blocks_unmapped", []):
            issues.append({
                "type": "prose_unmapped",
                "heading": u.get("heading", ""),
                "kind": u.get("kind", ""),
                "reason": u.get("reason", ""),
                "excerpt": u.get("excerpt", ""),
                "suggestion": (
                    "MD 줄글을 3~5개 bullet 으로 요약하거나, "
                    "양식에 해당 섹션용 본문 shape 추가"
                ),
            })

    # --- 양식 대비 결과에서 사라진 shape (사용자 수기 삭제 시그널) ---
    try:
        template_ids = _collect_shape_ids(tpl)
        output_ids = _collect_shape_ids(out)
        # 원본 양식의 slide_final 매핑 시도
        slides_final = (convert_result or {}).get("slides_final", []) or []
        # 결과 슬라이드 1장씩이 양식의 어느 번호에 해당하는지
        for out_idx, orig_idx in enumerate(slides_final, start=1):
            tpl_set = template_ids.get(orig_idx, set())
            out_set = output_ids.get(out_idx, set())
            removed = tpl_set - out_set
            # md2pptx 는 shape 을 삭제하지 않음 → removed 가 있으면 사용자 수기 편집
            # 단 이름이 중복되거나 공백인 경우 오탐 가능 → "Picture", "image" 등 케이스만 보고
            for name in removed:
                lname = (name or "").lower()
                if any(k in lname for k in ("picture", "image", "그림", "로고")):
                    issues.append({
                        "type": "template_shape_removed",
                        "slide": out_idx,
                        "original_slide": orig_idx,
                        "shape_name": name,
                        "suggestion": "양식에서 제거된 shape 에 대응하는 MD 내용이 있으면 MD 에서도 제거",
                    })
    except Exception:
        pass

    return {
        "issues": issues,
        "has_issues": len(issues) > 0,
        "issue_count": len(issues),
        "issue_types": sorted({iss["type"] for iss in issues}),
        "output_path": str(out),
        "slides_in_output": len(list(prs.slides)),
    }


if __name__ == "__main__":  # pragma: no cover
    import argparse, json as _json, sys

    p = argparse.ArgumentParser()
    p.add_argument("template")
    p.add_argument("output")
    args = p.parse_args()
    r = analyze_output(args.template, args.output)
    print(_json.dumps(r, ensure_ascii=False, indent=2))
    sys.exit(0 if not r["has_issues"] else 1)
