"""
PPTX 템플릿 조작: python-pptx로 슬라이드 텍스트/표 셀만 교체하고 나머지는 유지.

HWPX 버전 (doc_mcp/hwpx_vision/lib/hwpx_template.py) 과 철학 동일:
- 레이아웃/이미지/차트/애니메이션 건드리지 않음
- 제목 placeholder + 본문 placeholder + 표 셀 텍스트만 교체
- 폰트/크기/색상은 원본 run property 보존 (text 만 바꿈)
"""
from __future__ import annotations

import copy
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation


@dataclass
class SlideMeta:
    index: int
    title: str
    body_shapes: int = 0
    table_rows_cols: list[tuple[int, int]] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            "index": self.index,
            "title": self.title,
            "body_shapes": self.body_shapes,
            "has_table": bool(self.table_rows_cols),
            "table_rows_cols": self.table_rows_cols,
        }


def _extract_slide_title(slide) -> str:
    """슬라이드 제목 placeholder 텍스트 반환 (없으면 첫 텍스트 도형 사용)."""
    try:
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            return slide.shapes.title.text_frame.text.strip()
    except Exception:
        pass
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip().split("\n")[0]
    return ""


def _is_title_shape(shape) -> bool:
    """placeholder_format.idx==0 (title) 또는 ctrTitle(idx==?) 검사."""
    try:
        if getattr(shape, "is_placeholder", False):
            ph = shape.placeholder_format
            if ph.idx == 0:
                return True
            t = str(getattr(ph, "type", ""))
            if "TITLE" in t.upper():
                return True
    except Exception:
        pass
    return False


def _iter_body_shapes(slide):
    """제목이 아닌 텍스트 도형들을 순회."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if _is_title_shape(shape):
            continue
        yield shape


def list_slides(template_pptx: str) -> list[dict]:
    """슬라이드별 메타 정보 반환."""
    prs = Presentation(template_pptx)
    out: list[dict] = []
    for idx, slide in enumerate(prs.slides):
        meta = SlideMeta(index=idx, title=_extract_slide_title(slide))
        for shape in _iter_body_shapes(slide):
            meta.body_shapes += 1
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                meta.table_rows_cols.append((len(tbl.rows), len(tbl.columns)))
        out.append(meta.to_dict())
    return out


def _set_text_preserving_style(text_frame, new_text: str) -> None:
    """
    text_frame 의 텍스트를 new_text 로 교체하되, 첫 paragraph/run 의 서식을 재사용.
    줄바꿈(\n)은 새 paragraph 로 분리.
    """
    if text_frame is None:
        return
    paragraphs = text_frame.paragraphs
    if not paragraphs:
        text_frame.text = new_text
        return

    tpl_para = paragraphs[0]
    # 기존 모든 paragraph 요소 비우기 위해 XML 수준에서 접근
    p_elem = tpl_para._p
    txBody = p_elem.getparent()

    for p in list(txBody.findall(".//" + p_elem.tag)):
        txBody.remove(p)

    lines = new_text.split("\n") if new_text else [""]
    for i, line in enumerate(lines):
        new_p = copy.deepcopy(p_elem)
        # run 하나만 유지
        ns = p_elem.nsmap.get("a") or "http://schemas.openxmlformats.org/drawingml/2006/main"
        runs = new_p.findall(f"{{{ns}}}r")
        if runs:
            first_r = runs[0]
            for extra in runs[1:]:
                new_p.remove(extra)
            t_elem = first_r.find(f"{{{ns}}}t")
            if t_elem is not None:
                t_elem.text = line
        else:
            text_frame.add_paragraph().text = line
            continue
        # 기타 텍스트 관련 자식 제거 (br 등)
        for br in new_p.findall(f"{{{ns}}}br"):
            new_p.remove(br)
        txBody.append(new_p)


def _fill_table(tbl, md_body: str) -> bool:
    """
    LLM 본문을 '|' 구분 행으로 파싱해 표의 본문 셀에 주입.
    - 헤더(첫 행) 유지
    - 본문 행: 파싱된 행 수만큼 순차 대체, 열 개수 불일치 시 공백으로 채움
    - 파싱된 행이 없거나 적으면 남은 원본 행 유지 (덮어쓰지 않음)
    반환: 뭔가 셀을 교체했는지 여부
    """
    lines = [ln.strip() for ln in (md_body or "").split("\n") if ln.strip() and "|" in ln]
    if not lines:
        return False
    rows = []
    for ln in lines:
        parts = [p.strip() for p in ln.strip("|").split("|")]
        if parts:
            rows.append(parts)
    if not rows:
        return False

    cols = len(tbl.columns)
    body_row_start = 1  # 헤더 유지
    filled = False
    for r_offset, row_cells in enumerate(rows):
        r_idx = body_row_start + r_offset
        if r_idx >= len(tbl.rows):
            break
        for c_idx in range(cols):
            cell = tbl.cell(r_idx, c_idx)
            value = row_cells[c_idx] if c_idx < len(row_cells) else ""
            _set_text_preserving_style(cell.text_frame, value)
            filled = True
    return filled


def _norm(s: str) -> str:
    import re

    return re.sub(r"[\s.,:;~\-()·]+", "", s).lower()


def _match_slide(slide_title: str, candidates: dict[str, str]) -> Optional[str]:
    if slide_title in candidates:
        return slide_title
    st = _norm(slide_title)
    for k in candidates:
        if _norm(k) == st:
            return k
    for k in candidates:
        if slide_title and slide_title in k:
            return k
        if k and k in slide_title:
            return k
    return None


def inject_to_template(
    template_pptx: str,
    slide_to_body: dict[str, str],
    output_pptx: str,
) -> dict:
    """
    slide_to_body: {slide_title: body_text}
    각 슬라이드의 제목 텍스트로 매칭 → body_text 로 본문/표 교체.
    반환: {path, bytes, slides_replaced, matched_titles}
    """
    prs = Presentation(template_pptx)
    replaced = 0
    matched_titles: list[str] = []

    for slide in prs.slides:
        title = _extract_slide_title(slide)
        key = _match_slide(title, slide_to_body)
        if key is None:
            continue
        body_text = slide_to_body[key]

        table_written = False
        for shape in slide.shapes:
            if shape.has_table and not table_written:
                if _fill_table(shape.table, body_text):
                    table_written = True

        if not table_written:
            body_shapes = list(_iter_body_shapes(slide))
            if body_shapes:
                # 첫 본문 도형에 전체 본문 주입
                _set_text_preserving_style(body_shapes[0].text_frame, body_text)

        replaced += 1
        matched_titles.append(title)

    Path(output_pptx).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    size = Path(output_pptx).stat().st_size
    return {
        "path": output_pptx,
        "bytes": size,
        "slides_replaced": replaced,
        "matched_titles": matched_titles,
    }
